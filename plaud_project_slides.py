#!/usr/bin/env python3
"""
Plaud Project Summary - 技术面试 PPT
Multi-Agent 项目级摘要服务架构与实现
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ═══════════════════════════════════════════════════════════
# Design System Constants
# ═══════════════════════════════════════════════════════════

PARCHMENT   = RGBColor(0xf5, 0xf4, 0xed)
IVORY       = RGBColor(0xfa, 0xf9, 0xf5)
BRAND       = RGBColor(0x1B, 0x36, 0x5D)
BRAND_DEEP  = RGBColor(0x1B, 0x36, 0x5D)
NEAR_BLACK  = RGBColor(0x14, 0x14, 0x13)
DARK_WARM   = RGBColor(0x3d, 0x3d, 0x3a)
CHARCOAL    = RGBColor(0x4d, 0x4c, 0x48)
OLIVE       = RGBColor(0x50, 0x4e, 0x49)
STONE       = RGBColor(0x6b, 0x6a, 0x64)
BORDER      = RGBColor(0xe8, 0xe6, 0xdc)
WHITE       = RGBColor(0xff, 0xff, 0xff)

CN_SERIF = "Source Han Serif SC"
SERIF = CN_SERIF
SANS  = SERIF

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ═══════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════

def blank_slide(prs, bg_color=PARCHMENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide

def add_text(slide, text, left, top, width, height,
             font=SANS, size=18, bold=False, italic=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT,
             vanchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = vanchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_line(slide, left, top, width, color=BRAND, weight_pt=1):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight_pt)
    return line

def add_card(slide, left, top, width, height,
             fill=IVORY, border=BORDER, border_weight=0.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = border
    card.line.width = Pt(border_weight)
    card.shadow.inherit = False
    return card

# ═══════════════════════════════════════════════════════════
# Slide Templates
# ═══════════════════════════════════════════════════════════

def cover_slide(prs, title, subtitle, author, date):
    s = blank_slide(prs)
    add_text(s, title, Inches(1), Inches(2.5), Inches(11.33), Inches(1.5),
             font=SERIF, size=44, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.2), Inches(1), weight_pt=1.5)
    add_text(s, subtitle, Inches(1), Inches(4.5), Inches(11.33), Inches(0.8),
             font=SANS, size=18, color=OLIVE, align=PP_ALIGN.CENTER)
    add_text(s, f"{author}　·　{date}", Inches(1), Inches(6.5), Inches(11.33), Inches(0.4),
             font=SANS, size=13, color=STONE, align=PP_ALIGN.CENTER)
    return s

def toc_slide(prs, items):
    s = blank_slide(prs)
    add_text(s, "目录", Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
             font=SERIF, size=32, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(1.8), Inches(11), weight_pt=1)
    for i, item in enumerate(items):
        y = Inches(2.4 + i * 0.9)
        add_text(s, f"0{i+1}", Inches(1.2), y, Inches(1), Inches(0.6),
                 font=SERIF, size=28, color=BRAND)
        add_text(s, item, Inches(2.4), y, Inches(9), Inches(0.6),
                 font=SERIF, size=22, color=NEAR_BLACK, vanchor=MSO_ANCHOR.MIDDLE)
    return s

def chapter_slide(prs, number, title):
    s = blank_slide(prs, bg_color=BRAND)
    add_text(s, f"0{number}", Inches(0.8), Inches(0.5), Inches(2), Inches(0.8),
             font=SERIF, size=26, color=WHITE)
    add_text(s, title, Inches(1), Inches(3), Inches(11.33), Inches(1.5),
             font=SERIF, size=56, color=WHITE, align=PP_ALIGN.CENTER)
    return s

def content_slide(prs, eyebrow, title, body, page_num=None):
    s = blank_slide(prs)
    add_text(s, eyebrow, Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=12, color=STONE)
    add_text(s, title, Inches(1.2), Inches(1.2), Inches(11.33), Inches(1.2),
             font=SERIF, size=32, color=NEAR_BLACK)
    add_text(s, body, Inches(1.2), Inches(3), Inches(11), Inches(3.5),
             font=SANS, size=18, color=DARK_WARM)
    if page_num is not None:
        add_text(s, f" - {page_num:02d}", Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE, align=PP_ALIGN.RIGHT)
    return s

def metrics_slide(prs, title, metrics):
    s = blank_slide(prs)
    add_text(s, title, Inches(1.2), Inches(0.8), Inches(11), Inches(1),
             font=SERIF, size=28, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(2), Inches(1))

    n = len(metrics)
    card_w = Inches(2.8)
    gap = Inches(0.3)
    total_w = card_w * n + gap * (n - 1)
    start = (SLIDE_W - total_w) / 2

    for i, (value, label) in enumerate(metrics):
        x = start + (card_w + gap) * i
        add_text(s, value, x, Inches(3), card_w, Inches(1.5),
                 font=SERIF, size=52, color=BRAND, align=PP_ALIGN.CENTER)
        add_text(s, label, x, Inches(4.8), card_w, Inches(0.6),
                 font=SANS, size=14, color=OLIVE, align=PP_ALIGN.CENTER)
    return s

def comparison_slide(prs, eyebrow, left_title, left_items, right_title, right_items, page_num=None):
    s = blank_slide(prs)
    add_text(s, eyebrow, Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=12, color=STONE)
    divider = s.shapes.add_connector(1, Inches(6.67), Inches(1.0), Inches(6.67), Inches(6.8))
    divider.line.color.rgb = BORDER
    divider.line.width = Pt(1)
    add_text(s, left_title, Inches(1.2), Inches(1.2), Inches(5), Inches(0.8),
             font=SERIF, size=22, color=OLIVE)
    add_text(s, right_title, Inches(7.0), Inches(1.2), Inches(5), Inches(0.8),
             font=SERIF, size=22, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(2.2), Inches(11.5), weight_pt=0.5)
    for i, item in enumerate(left_items[:4]):
        add_text(s, item, Inches(1.2), Inches(2.6 + i * 0.9), Inches(4.9), Inches(0.7),
                 font=SANS, size=17, color=STONE)
    for i, item in enumerate(right_items[:4]):
        add_text(s, item, Inches(7.0), Inches(2.6 + i * 0.9), Inches(5.2), Inches(0.7),
                 font=SANS, size=17, color=DARK_WARM)
    if page_num is not None:
        add_text(s, f" - {page_num:02d}", Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE, align=PP_ALIGN.RIGHT)
    return s

def pipeline_slide(prs, eyebrow, title, steps, page_num=None):
    s = blank_slide(prs)
    add_text(s, eyebrow, Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=12, color=STONE)
    add_text(s, title, Inches(1.2), Inches(1.1), Inches(11), Inches(0.9),
             font=SERIF, size=30, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(2.15), Inches(11), weight_pt=0.5)

    n = len(steps[:4])
    step_w = Inches(11.5 / n)
    for i, (step_title, step_desc) in enumerate(steps[:4]):
        x = Inches(1.0) + step_w * i
        add_text(s, f"0{i+1}", x, Inches(2.5), step_w, Inches(0.8),
                 font=SERIF, size=40, color=BRAND)
        add_text(s, step_title, x, Inches(3.45), step_w - Inches(0.2), Inches(0.6),
                 font=SERIF, size=19, color=NEAR_BLACK)
        add_text(s, step_desc, x, Inches(4.15), step_w - Inches(0.2), Inches(2.2),
                 font=SANS, size=15, color=OLIVE)
    if page_num is not None:
        add_text(s, f" - {page_num:02d}", Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE, align=PP_ALIGN.RIGHT)
    return s

def ending_slide(prs, message, contact):
    s = blank_slide(prs)
    add_text(s, message, Inches(1), Inches(3), Inches(11.33), Inches(1.2),
             font=SERIF, size=40, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.5), Inches(1), weight_pt=1.5)
    add_text(s, contact, Inches(1), Inches(4.8), Inches(11.33), Inches(0.6),
             font=SANS, size=16, color=OLIVE, align=PP_ALIGN.CENTER)
    return s

def bullet_slide(prs, eyebrow, title, bullets, page_num=None):
    """带编号列表的内容页"""
    s = blank_slide(prs)
    add_text(s, eyebrow, Inches(1.2), Inches(0.6), Inches(10), Inches(0.4),
             font=SANS, size=12, color=STONE)
    add_text(s, title, Inches(1.2), Inches(1.2), Inches(11), Inches(1),
             font=SERIF, size=32, color=NEAR_BLACK)
    add_line(s, Inches(1.2), Inches(2.35), Inches(11), weight_pt=0.5)

    for i, bullet in enumerate(bullets[:5]):
        y = Inches(2.8 + i * 0.8)
        add_text(s, f"• {bullet}", Inches(1.4), y, Inches(10.8), Inches(0.7),
                 font=SANS, size=17, color=DARK_WARM)
    if page_num is not None:
        add_text(s, f" - {page_num:02d}", Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.3),
                 font=SANS, size=11, color=STONE, align=PP_ALIGN.RIGHT)
    return s

# ═══════════════════════════════════════════════════════════
# Main Content
# ═══════════════════════════════════════════════════════════

def main():
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default="plaud_project_summary.pptx")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    page = 1

    # 01 封面
    cover_slide(prs,
        title="Plaud Project Summary",
        subtitle="Multi-Agent 项目级摘要服务 · 架构与实现",
        author="梁柱",
        date="2026.06")

    # 02 目录
    toc_slide(prs, items=[
        "项目背景与核心价值",
        "整体架构生命周期",
        "Agent 工作流详解",
        "关键技术决策",
    ])

    # ===== 第 1 章：项目背景 =====
    chapter_slide(prs, 1, "项目背景与核心价值")

    # 03 电梯陈述
    page += 1
    content_slide(prs,
        eyebrow="01 · 项目定位",
        title="从 0 到 1 独立负责的 Multi-Agent 摘要服务",
        body="把一个项目下的多份录音和文档聚合起来，产出结构化项目报告\n\n支持三种策略：总结、对比、进度跟踪\n\n核心是基于 DeepAgents 的 Agent 工作流，配合 Temporal 异步编排、引用溯源工程和自研评测体系",
        page_num=page)

    # 04 要解决的问题
    page += 1
    bullet_slide(prs,
        eyebrow="01 · 核心挑战",
        title="单次摘要无法回答跨文件问题",
        bullets=[
            "用户真实场景：一个项目持续数周，积累十几份录音和文档",
            "需要跨文件才能回答的问题：对比候选人、追踪项目进展",
            "三种策略对应三类诉求：SUMMARY 合并去重、COMPARISON 横向对比、PROGRESS 跟踪趋势"
        ],
        page_num=page)

    # 05 技术挑战
    page += 1
    pipeline_slide(prs,
        eyebrow="01 · 技术难点",
        title="端到端上线的三大技术挑战",
        steps=[
            ("Context 不可控", "少则两三份、多则十几份，token 量跨度极大，单一处理方式扛不住"),
            ("长耗时任务", "Multi-agent 流程几十秒级，同步 HTTP 必然超时"),
            ("可信与可量化", "报告必须能溯源到原文件，改一版 prompt 得能客观判断好坏"),
        ],
        page_num=page)

    # 06 我的角色
    page += 1
    metrics_slide(prs,
        title="独立负责，端到端 0→1",
        metrics=[
            ("1 人", "需求拆解到上线运维"),
            ("3 种", "摘要策略"),
            ("生产", "已上线管道"),
        ])

    # ===== 第 2 章：整体架构 =====
    chapter_slide(prs, 2, "整体架构生命周期")

    # 07 请求生命周期
    page += 1
    pipeline_slide(prs,
        eyebrow="02 · 请求流",
        title="一个请求的完整生命周期",
        steps=[
            ("Temporal 触发", "上游 Workflow 发 HTTP POST"),
            ("立即返回 202", "API 层校验后丢进后台 asyncio.Task"),
            ("Agent 异步执行", "预处理 → 语言检测 → Agent → 提取标题 → 存 S3"),
            ("Signal 回传", "结果通过 Temporal Signal 异步回传上游"),
        ],
        page_num=page)

    # 08 Temporal 外部服务模式
    page += 1
    content_slide(prs,
        eyebrow="02 · 架构模式",
        title="Temporal 外部服务模式解耦长耗时任务",
        body="上游 Workflow 发 HTTP 触发，本服务立即返回 202\n\n算完用 Temporal Signal 把结果异步回传给那个 Workflow\n\nHTTP 只负责收下任务，结果走 Signal，彻底解耦",
        page_num=page)

    # 09 K8s 优雅关停
    page += 1
    pipeline_slide(prs,
        eyebrow="02 · 优雅关停",
        title="K8s 滚动发布时保护在途任务",
        steps=[
            ("收到 SIGTERM", "置 _shutting_down 标志"),
            ("新请求返 503", "健康检查也返 503，K8s 停止转发流量"),
            ("等待任务完成", "asyncio.gather 追踪 _running_tasks"),
            ("安全退出", "保证在途长任务不被腰斩"),
        ],
        page_num=page)

    # ===== 第 3 章：Agent 工作流 =====
    chapter_slide(prs, 3, "Agent 工作流详解")

    # 10 编排范式
    page += 1
    content_slide(prs,
        eyebrow="03 · 编排哲学",
        title="引导式工具编排 + 硬约束兜底",
        body="System prompt 把流程写成 5 步顺序，引导模型按序调工具\n\n工具限流给每个工具设调用上限，防止死循环和乱序\n\n能用规则保证的就别指望 prompt，能用 prompt 表达的就别写死成代码",
        page_num=page)

    # 11 文件级摘要预处理
    page += 1
    content_slide(prs,
        eyebrow="03 · 预处理",
        title="先给每个文件生成 file_summary，压低后续 context",
        body="跑 Agent 之前，API 层先调 ensure_file_summaries\n\n每个文件生成一份文件级摘要，同时做语言检测\n\n摘要来源按固定优先级降级：已有摘要 → Filesystem → Temporal Workflow → 本地 LLM",
        page_num=page)

    # 12 SummaryState 数据总线
    page += 1
    bullet_slide(prs,
        eyebrow="03 · 状态设计",
        title="SummaryState 是工具之间唯一的数据总线",
        bullets=[
            "工具之间从不互相调用，全部通过共享 SummaryState 读写传递数据",
            "上一个工具把结果写进 state 字段，下一个工具从里面读",
            "Map/Reduce 并发写用 reducer 合并，防止覆盖",
            "主从 Agent 共享同一份 state，消息隔离",
        ],
        page_num=page)

    # 13 路由决策
    page += 1
    comparison_slide(prs,
        eyebrow="03 · 路由",
        left_title="ONE-PASS（小内容）",
        left_items=[
            "所有内容一次性喂给 LLM",
            "简单、快、上下文连贯",
            "默认阈值约 60K tokens"
        ],
        right_title="MAP-REDUCE（大内容）",
        right_items=[
            "拆解后分而治之",
            "绕开单次 context window 限制",
            "避免「中间信息丢失」问题"
        ],
        page_num=page)

    # 14 共有规划阶段
    page += 1
    pipeline_slide(prs,
        eyebrow="03 · 规划",
        title="两条路都先走三个规划工具",
        steps=[
            ("分析文件关系", "LLM 分析文件间关系、识别分歧，结果跨任务复用"),
            ("生成大纲 + RAG 查询", "LLM 生成文档大纲，对 location_file 用 structured output 抽 RAG query"),
            ("RAG 检索补全", "对 location_file 走 Retriever API 检索，回填 rag_content"),
        ],
        page_num=page)

    # 15 ONE-PASS 路径
    page += 1
    content_slide(prs,
        eyebrow="03 · ONE-PASS",
        title="大纲 + 所有文件内容一次性交给 LLM",
        body="generate_few_file_summary 工具直接生成 final_summary\n\n内容能塞进一个 context window 时，这条路上下文最连贯、调用最少",
        page_num=page)

    # 16 MAP-REDUCE 路径
    page += 1
    pipeline_slide(prs,
        eyebrow="03 · MAP-REDUCE",
        title="委托给独立子 Agent，三步流程",
        steps=[
            ("Map", "每个文件并行抽取 section，按文件切"),
            ("Reduce", "换维度：把所有文件的同一个 section 并行合并，按 section 切"),
            ("Assembly", "按大纲组装、补过渡语，得到 final_summary"),
        ],
        page_num=page)

    # 17 主从 Agent 通信
    page += 1
    bullet_slide(prs,
        eyebrow="03 · 子 Agent",
        title="主从 Agent 的 state 共享与 context 隔离",
        bullets=[
            "SubAgentMiddleware 把子 Agent 暴露成一个 task 工具",
            "主 Agent 用 task(subagent_type='map-reduce-synthesizer') 发一次 tool call",
            "主从共享同一份 SummaryState，子 Agent 直接读写",
            "子 Agent 内部几十轮消息不回流主 Agent，只返回一条完成状态",
            "主 Agent 的 context 始终干净，海量中间产物被隔离在子 Agent 里"
        ],
        page_num=page)

    # 18 收尾 review
    page += 1
    content_slide(prs,
        eyebrow="03 · 收尾",
        title="review_final_summary 终审润色",
        body="回到主 Agent，对 final_summary 做一遍终审润色\n\n检测文档类型、调整风格、保留引用标记\n\nSystem prompt 明确要求这一步做完工作流即结束、不再调任何工具",
        page_num=page)

    # 19 健壮性两道防线
    page += 1
    comparison_slide(prs,
        eyebrow="03 · 健壮性",
        left_title="Middleware 栈",
        left_items=[
            "AgentLoggingMiddleware 执行日志",
            "TodoListMiddleware 任务追踪",
            "SummarizationMiddleware 170K token 压缩",
            "PatchToolCallsMiddleware 修正格式错误"
        ],
        right_title="工具限流 + 锁写",
        right_items=[
            "每个工具调用上限（生成类 2 次、规划类 1 次）",
            "review 后锁写，拒绝再执行生成类工具",
            "防止 Agent 又跑回去重新生成覆盖润色结果",
            "防止死循环和反复调同一工具"
        ],
        page_num=page)

    # ===== 第 4 章：关键决策 =====
    chapter_slide(prs, 4, "关键技术决策")

    # 20 Citation 引用溯源
    page += 1
    pipeline_slide(prs,
        eyebrow="04 · 引用工程",
        title="LLM 内部用短 ID，出口做确定性还原",
        steps=[
            ("LLM 写短 ID", "省 token，<<cite:source_id>> 格式"),
            ("三级匹配还原", "精确 → 后缀 → 子串，每级只在唯一命中时才还原"),
            ("替换回文件名", "正文里裸露的 source_id 替换回可读文件名"),
            ("不再过 LLM", "引用还原是确定性匹配问题，交给 LLM 只会引入幻觉"),
        ],
        page_num=page)

    # 21 自研评测体系
    page += 1
    bullet_slide(prs,
        eyebrow="04 · 评测",
        title="G-Eval 模式：多次采样 + 加权维度",
        bullets=[
            "每种策略配一套加权维度（SUMMARY 有 7 个维度，去重合并质量权重 1.5）",
            "每个维度多次采样（N=8，temperature 0.7）取均值，降低随机性",
            "8 个维度并行评测（ThreadPoolExecutor），Judge 用强模型",
            "每次改动都能跑出可量化、可对比的维度分和总分",
        ],
        page_num=page)

    # 22 区域感知路由
    page += 1
    content_slide(prs,
        eyebrow="04 · 多区域",
        title="按 AWS_REGION 自动选模型与框架",
        body="海外用 Gemini（Vertex AI），国内用豆包（Model Hub SDK）\n\n凭证、路由、prompt、各工具的模型配置托管在 AWS AppConfig\n\nWatchdog 线程轮询（120s）检测变更，多数配置零停机热更新",
        page_num=page)

    # 23 技术栈总结
    page += 1
    bullet_slide(prs,
        eyebrow="04 · 技术栈",
        title="端到端覆盖 Agent / 后端 / 编排 / 运维",
        bullets=[
            "Agent / LLM: DeepAgents、LangChain、Langfuse、Gemini、豆包、GPT（Judge）",
            "后端: FastAPI、asyncio 后台任务",
            "编排: Temporal（外部服务模式 + Signal 回传）",
            "存储: AWS S3 双桶、Filesystem KV、Retriever API（RAG）",
            "运维: K8s、Supervisor、AWS AppConfig、OpenTelemetry"
        ],
        page_num=page)

    # 24 结束
    ending_slide(prs,
        message="Thank you",
        contact="Q & A")

    prs.save(args.out)
    print(f"✓ 已生成: {args.out}")
    print(f"  共 {len(prs.slides)} 张幻灯片")

if __name__ == '__main__':
    main()
