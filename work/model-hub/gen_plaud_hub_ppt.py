#!/usr/bin/env python3
"""生成 Plaud Model Hub 技术分享 PPT。

用法:
  pip install python-pptx
  python3 gen_plaud_hub_ppt.py --out PlaudModelHub技术分享.pptx

样式参考 kami parchment 设计系统:
- 暖纸底色 #f5f4ed
- 单一墨蓝强调色 #1B365D
- 一律衬线字体 Source Han Serif SC
- 16:9 宽屏 13.33 × 7.5 inch
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ═══════════════════════════════════════════════════════════
# 设计系统常量
# ═══════════════════════════════════════════════════════════

PARCHMENT   = RGBColor(0xf5, 0xf4, 0xed)
IVORY       = RGBColor(0xfa, 0xf9, 0xf5)
BRAND       = RGBColor(0x1B, 0x36, 0x5D)
BRAND_TINT  = RGBColor(0xEE, 0xF2, 0xF7)
NEAR_BLACK  = RGBColor(0x14, 0x14, 0x13)
DARK_WARM   = RGBColor(0x3d, 0x3d, 0x3a)
OLIVE       = RGBColor(0x50, 0x4e, 0x49)
STONE       = RGBColor(0x6b, 0x6a, 0x64)
BORDER      = RGBColor(0xe8, 0xe6, 0xdc)
WHITE       = RGBColor(0xff, 0xff, 0xff)

SERIF = "Source Han Serif SC"
SANS  = SERIF
MONO  = "JetBrains Mono"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ═══════════════════════════════════════════════════════════
# 基础原语
# ═══════════════════════════════════════════════════════════

def blank_slide(prs, bg=PARCHMENT):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    rect.shadow.inherit = False
    return s


def add_text(slide, text, left, top, width, height,
             font=SANS, size=14, bold=False,
             color=NEAR_BLACK, align=PP_ALIGN.LEFT,
             vanchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = vanchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_multi_text(slide, lines, left, top, width, height,
                   font=SANS, size=14, color=DARK_WARM,
                   line_spacing=1.4):
    """多行文本,每行可单独配 (text, color, bold)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        if isinstance(item, tuple):
            text, *rest = item
            c = rest[0] if len(rest) > 0 else color
            b = rest[1] if len(rest) > 1 else False
            sz = rest[2] if len(rest) > 2 else size
        else:
            text, c, b, sz = item, color, False, size
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(sz)
        run.font.bold = b
        run.font.color.rgb = c
    return tb


def add_rect(slide, left, top, width, height,
             fill=IVORY, border=BORDER, border_pt=0.5,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if border is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = border
        sh.line.width = Pt(border_pt)
    sh.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sh.adjustments[0] = 0.08
        except Exception:
            pass
    return sh


def add_line(slide, x1, y1, x2, y2, color=BORDER, pt=0.5):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(pt)
    return ln


def add_arrow(slide, x1, y1, x2, y2, color=STONE, pt=1):
    from pptx.oxml.ns import qn
    from lxml import etree
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(pt)
    # 加箭头头部
    line_xml = ln.line._get_or_add_ln()
    tail = etree.SubElement(line_xml, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'sm')
    tail.set('len', 'sm')
    return ln


def add_eyebrow(slide, text, page_num=None):
    add_text(slide, text,
             Inches(0.6), Inches(0.4), Inches(10), Inches(0.4),
             font=MONO, size=10, color=STONE)
    if page_num is not None:
        add_text(slide, f"{page_num:02d} / 19",
                 Inches(12), Inches(0.4), Inches(1), Inches(0.4),
                 font=MONO, size=10, color=STONE, align=PP_ALIGN.RIGHT)


def add_title(slide, text, top=0.9, size=26):
    add_text(slide, text,
             Inches(0.6), Inches(top), Inches(12.13), Inches(1.2),
             font=SERIF, size=size, color=NEAR_BLACK, bold=True)


def add_lead(slide, text, top=1.85):
    add_text(slide, text,
             Inches(0.6), Inches(top), Inches(12.13), Inches(0.6),
             font=SANS, size=14, color=OLIVE)


def add_footer_note(slide, text):
    add_rect(slide, Inches(0.6), Inches(6.85), Inches(12.13), Inches(0.4),
             fill=BRAND_TINT, border=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, text,
             Inches(0.8), Inches(6.88), Inches(11.7), Inches(0.34),
             font=SANS, size=12, color=BRAND, bold=True,
             vanchor=MSO_ANCHOR.MIDDLE)


# ═══════════════════════════════════════════════════════════
# slide 1 · 封面
# ═══════════════════════════════════════════════════════════

def slide_cover(prs):
    s = blank_slide(prs)
    add_text(s, "Plaud Model Hub",
             Inches(0.6), Inches(2.4), Inches(12.13), Inches(1.0),
             font=SERIF, size=44, color=NEAR_BLACK, align=PP_ALIGN.CENTER)
    add_text(s, "多供应商 LLM 调度的统一抽象",
             Inches(0.6), Inches(3.4), Inches(12.13), Inches(1.0),
             font=SERIF, size=36, color=BRAND, align=PP_ALIGN.CENTER, bold=True)
    add_line(s, Inches(6.17), Inches(4.7), Inches(7.17), Inches(4.7),
             color=BRAND, pt=1.5)
    add_text(s, "从三种调用模式到自适应权重 · 一次讲透架构与容错",
             Inches(0.6), Inches(4.85), Inches(12.13), Inches(0.5),
             font=SANS, size=16, color=OLIVE, align=PP_ALIGN.CENTER)
    add_text(s, "技术分享 · 2026",
             Inches(0.6), Inches(6.6), Inches(12.13), Inches(0.4),
             font=MONO, size=11, color=STONE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# slide 2 · 目录
# ═══════════════════════════════════════════════════════════

def slide_toc(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "AGENDA · 18 章")
    add_title(s, "目录")

    items_left = [
        ("01", "痛点 · 多供应商让业务侧重复造轮子"),
        ("02", "三种调用模式 · 统一抽象的核心选择"),
        ("03", "整体架构地图 · 五个关键层"),
        ("04", "ModelRequest · 透传四字段"),
        ("05", "CoreEngine · 循环外/循环内"),
        ("06", "★ Provider 三岔分派 · 核心抽象"),
        ("07", "三层容错总览 · 两条独立轨道"),
        ("08", "救援链 · Retry → Failover → Fallback"),
        ("09", "熔断器 · 与自适应权重的分工"),
    ]
    items_right = [
        ("10", "★ 自适应权重 · 429 是太忙不是坏了"),
        ("11", "★ 权重重分配 · 必要性"),
        ("12", "权重重分配 · 公式"),
        ("13", "权重重分配 · 完整数值演练"),
        ("14", "★ recovery_rate · 限速恢复"),
        ("15", "完整调用旅程 · 跨厂商 fallback"),
        ("16", "能力复用矩阵 · 三模式共享"),
        ("17", "设计哲学 · 四条原则"),
        ("18", "总结"),
    ]

    for i, (num, title) in enumerate(items_left):
        y = Inches(1.95 + i * 0.5)
        add_text(s, num, Inches(0.8), y, Inches(0.7), Inches(0.4),
                 font=SERIF, size=18, color=BRAND, bold=True)
        add_text(s, title, Inches(1.55), y + Inches(0.03), Inches(5.0), Inches(0.4),
                 font=SANS, size=14, color=NEAR_BLACK)

    for i, (num, title) in enumerate(items_right):
        y = Inches(1.95 + i * 0.5)
        add_text(s, num, Inches(7.0), y, Inches(0.7), Inches(0.4),
                 font=SERIF, size=18, color=BRAND, bold=True)
        add_text(s, title, Inches(7.75), y + Inches(0.03), Inches(5.0), Inches(0.4),
                 font=SANS, size=14, color=NEAR_BLACK)


# ═══════════════════════════════════════════════════════════
# slide 3 · 痛点
# ═══════════════════════════════════════════════════════════

def slide_pain(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "01 · 出发点", 3)
    add_title(s, "同时接 8 家 LLM 厂商,重复造轮子的代价已经不可控")
    add_lead(s, "每家 SDK 不同、错误结构不同、配额限流差异巨大;让业务自己处理,每个团队都得写一遍 client、fallback、重试。")

    cards = [
        ("SDK 异构", "OpenAI · Anthropic · Gemini\n火山 · DashScope · Bedrock\n各自一套 client,错误格式互不兼容"),
        ("配额倒挂", "同一逻辑模型常买 2 至 3 个\nendpoint,一个限流业务直接挂"),
        ("切换成本", "换模型 = 改代码 + 发版\n做负载均衡得自己写"),
        ("容错缺失", "没有统一的 429 降权\n熔断与 fallback 链路"),
    ]

    card_w = Inches(2.93)
    gap = Inches(0.18)
    start = Inches(0.6)
    top = Inches(2.9)
    h = Inches(3.4)

    for i, (title, body) in enumerate(cards):
        x = start + (card_w + gap) * i
        add_rect(s, x, top, card_w, h, fill=IVORY, border=BORDER, border_pt=0.6)
        # 顶端品牌色短条
        add_rect(s, x + Inches(0.25), top + Inches(0.3), Inches(0.5), Inches(0.05),
                 fill=BRAND, border=None, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, title,
                 x + Inches(0.25), top + Inches(0.55), card_w - Inches(0.5), Inches(0.5),
                 font=SERIF, size=18, color=NEAR_BLACK, bold=True)
        add_text(s, body,
                 x + Inches(0.25), top + Inches(1.15), card_w - Inches(0.5), h - Inches(1.3),
                 font=SANS, size=13, color=DARK_WARM)

    add_footer_note(s, "基础设施层兜住共性,业务方只写 app_id 与 logical_model 两个名字。")


# ═══════════════════════════════════════════════════════════
# slide 4 · 三种调用模式
# ═══════════════════════════════════════════════════════════

def slide_modes(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "02 · 核心思路", 4)
    add_title(s, "把所有调用统一成一个信封,三种模式按透传与风格分流")
    add_lead(s, "CoreEngine 对外只有一个签名 invoke(request) → response,所有能力栈完整复用。")

    modes = [
        {
            "tag": "DISABLED",
            "name": "统一模式",
            "body": "业务方填 messages /\ntemperature 等统一字段",
            "cost": "2 次转换",
            "subcost": "hub 抽象 ↔ 各家 SDK",
            "suit": "适合不关心厂商差异的业务\n要的是统一抽象与切换自由",
        },
        {
            "tag": "SAME_STYLE",
            "name": "同风格透传",
            "body": "业务方用原生 SDK 写法\n路由到本家 endpoint",
            "cost": "0 次转换",
            "subcost": "原样喂 SDK,无损",
            "suit": "适合已有原生 SDK 代码\n不想改写的业务",
        },
        {
            "tag": "CROSS_STYLE",
            "name": "跨风格透传",
            "body": "业务方用原生 SDK 写法\n路由到异家 endpoint",
            "cost": "2 次转换",
            "subcost": "adapter 内部翻译",
            "suit": "适合想跨厂商兜底\n且不愿改 SDK 代码",
        },
    ]

    card_w = Inches(3.93)
    gap = Inches(0.22)
    start = Inches(0.6)
    top = Inches(2.7)
    h = Inches(3.95)

    for i, m in enumerate(modes):
        x = start + (card_w + gap) * i
        add_rect(s, x, top, card_w, h, fill=IVORY, border=BRAND, border_pt=1)
        # 标签
        add_rect(s, x + Inches(0.3), top + Inches(0.3), Inches(1.6), Inches(0.35),
                 fill=BRAND, border=None, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, m["tag"],
                 x + Inches(0.3), top + Inches(0.3), Inches(1.6), Inches(0.35),
                 font=MONO, size=10, color=WHITE, align=PP_ALIGN.CENTER,
                 vanchor=MSO_ANCHOR.MIDDLE, bold=True)
        # 模式名
        add_text(s, m["name"],
                 x + Inches(0.3), top + Inches(0.8), card_w - Inches(0.6), Inches(0.5),
                 font=SERIF, size=20, color=NEAR_BLACK, bold=True)
        # 调用形态
        add_text(s, m["body"],
                 x + Inches(0.3), top + Inches(1.45), card_w - Inches(0.6), Inches(0.85),
                 font=SANS, size=12, color=DARK_WARM)
        # 转换次数
        add_text(s, m["cost"],
                 x + Inches(0.3), top + Inches(2.4), card_w - Inches(0.6), Inches(0.4),
                 font=SERIF, size=18, color=BRAND, bold=True)
        add_text(s, m["subcost"],
                 x + Inches(0.3), top + Inches(2.78), card_w - Inches(0.6), Inches(0.3),
                 font=SANS, size=11, color=STONE)
        # 分割线
        add_line(s, x + Inches(0.3), top + Inches(3.2),
                 x + card_w - Inches(0.3), top + Inches(3.2), color=BORDER, pt=0.5)
        # 适合谁
        add_text(s, m["suit"],
                 x + Inches(0.3), top + Inches(3.3), card_w - Inches(0.6), Inches(0.6),
                 font=SANS, size=11, color=OLIVE)


# ═══════════════════════════════════════════════════════════
# slide 5 · 整体架构地图
# ═══════════════════════════════════════════════════════════

def slide_architecture(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "03 · 整体地图", 5)
    add_title(s, "请求从业务入口到原生 SDK 的旅程,只经过五个关键层")

    # 画布起点
    cx, cy = Inches(0.6), Inches(2.0)
    # 第一行: 业务入口 → ModelRequest → Transport → Engine
    row1_top = cy
    row1_h = Inches(1.2)
    boxes = [
        ("业务入口 · 三选一", "ChatModelHub · ModelHubClient\nWrapped OpenAI / Anthropic / GenAI", IVORY, BORDER),
        ("ModelRequest", "passthrough_mode + source_style\nraw_request + raw_method", BRAND_TINT, BRAND),
        ("Transport", "DirectProviderTransport\n纯方法转发,预留 HTTP 扩展点", IVORY, BORDER),
        ("CoreEngine · invoke()", "①validate ②before ③排除+权重\n④Router ⑤合并参数 ⑥call ⑦回写/救援", IVORY, BRAND),
    ]
    box_w = Inches(2.93)
    gap = Inches(0.16)
    for i, (title, body, fill, border) in enumerate(boxes):
        x = cx + (box_w + gap) * i
        bw = 1 if border == BRAND else 0.6
        add_rect(s, x, row1_top, box_w, row1_h, fill=fill, border=border, border_pt=bw)
        add_text(s, title, x + Inches(0.15), row1_top + Inches(0.12),
                 box_w - Inches(0.3), Inches(0.4),
                 font=SERIF, size=13, color=BRAND, bold=True)
        add_text(s, body, x + Inches(0.15), row1_top + Inches(0.5),
                 box_w - Inches(0.3), row1_h - Inches(0.55),
                 font=SANS, size=10, color=DARK_WARM)
        if i < 3:
            arr_x1 = x + box_w
            arr_x2 = arr_x1 + gap
            arr_y = row1_top + row1_h / 2
            add_arrow(s, arr_x1, arr_y, arr_x2, arr_y, color=STONE, pt=1)

    # 第二行: Provider 三岔
    row2_top = row1_top + row1_h + Inches(0.5)
    row2_h = Inches(1.5)
    total_w = box_w * 4 + gap * 3
    add_rect(s, cx, row2_top, total_w, row2_h,
             fill=BRAND_TINT, border=BRAND, border_pt=1)
    add_text(s, "ModelProvider._invoke · 三岔分派",
             cx + Inches(0.2), row2_top + Inches(0.12), total_w - Inches(0.4), Inches(0.4),
             font=SERIF, size=14, color=BRAND, bold=True, align=PP_ALIGN.CENTER)
    sub_w = (total_w - Inches(0.4) - gap * 2) / 3
    for i, (label, sub) in enumerate([
        ("_invoke_unified", "DISABLED · 2 次转换"),
        ("_invoke_passthrough", "SAME_STYLE · 0 次"),
        ("_invoke_adapted", "CROSS_STYLE · adapter"),
    ]):
        x = cx + Inches(0.2) + (sub_w + gap) * i
        add_rect(s, x, row2_top + Inches(0.6), sub_w, Inches(0.8),
                 fill=IVORY, border=BORDER, border_pt=0.5)
        add_text(s, label, x, row2_top + Inches(0.7), sub_w, Inches(0.3),
                 font=MONO, size=11, color=NEAR_BLACK, align=PP_ALIGN.CENTER, bold=True)
        add_text(s, sub, x, row2_top + Inches(1.05), sub_w, Inches(0.3),
                 font=SANS, size=10, color=OLIVE, align=PP_ALIGN.CENTER)

    # 箭头: Engine → Provider 三岔
    eng_x = cx + (box_w + gap) * 3 + box_w / 2
    add_arrow(s, eng_x, row1_top + row1_h, eng_x, row2_top, color=STONE, pt=1)

    # 第三行: ClientPool + 原生 SDK + ConfigStore
    row3_top = row2_top + row2_h + Inches(0.4)
    row3_h = Inches(0.85)
    for i, (title, body, fill) in enumerate([
        ("ConfigStore", "File · AppConfig · Http", IVORY),
        ("ClientPool", "按 endpoint_id 复用 SDK client,HTTP 连接保活", IVORY),
        ("原生 SDK", "openai · anthropic · google-genai · boto3", IVORY),
    ]):
        widths = [box_w, box_w * 1.5 + gap, box_w * 1.5 + gap]
        xs = [cx,
              cx + box_w + gap,
              cx + box_w + gap + box_w * 1.5 + gap + gap]
        x = xs[i]
        w = widths[i]
        add_rect(s, x, row3_top, w, row3_h, fill=fill, border=BORDER, border_pt=0.5)
        add_text(s, title, x + Inches(0.15), row3_top + Inches(0.1),
                 w - Inches(0.3), Inches(0.35),
                 font=SERIF, size=12, color=BRAND, bold=True)
        add_text(s, body, x + Inches(0.15), row3_top + Inches(0.45),
                 w - Inches(0.3), row3_h - Inches(0.5),
                 font=SANS, size=10, color=DARK_WARM)

    add_footer_note(s, "ConfigStore 通过 RuntimeConfig 向 Engine 注入逻辑模型、endpoint 配置与策略。")


# ═══════════════════════════════════════════════════════════
# slide 6 · ModelRequest 四字段
# ═══════════════════════════════════════════════════════════

def slide_modelrequest(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "04 · 信封字段", 6)
    add_title(s, "ModelRequest 的四个透传字段决定每一次请求的命运")

    # 左右两栏
    col_top = Inches(2.0)
    col_h = Inches(4.6)
    col_w = Inches(6.0)
    left_x = Inches(0.6)
    right_x = Inches(6.93)

    add_rect(s, left_x, col_top, col_w, col_h, fill=IVORY, border=BORDER, border_pt=0.5)
    add_rect(s, right_x, col_top, col_w, col_h, fill=IVORY, border=BORDER, border_pt=0.5)

    # 左侧 · 统一模式字段
    add_text(s, "统一模式字段",
             left_x + Inches(0.3), col_top + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             font=SERIF, size=18, color=BRAND, bold=True)
    add_multi_text(s, [
        ("logical_model", NEAR_BLACK, True, 13),
        ("    业务侧逻辑模型名,如 summary:gpt-4.1", OLIVE, False, 11),
        ("messages", NEAR_BLACK, True, 13),
        ("    统一消息结构", OLIVE, False, 11),
        ("temperature · max_tokens", NEAR_BLACK, True, 13),
        ("    顶层统一参数,Hub 接管翻译", OLIVE, False, 11),
        ("provider_params", NEAR_BLACK, True, 13),
        ("    Hub 不认识的厂商原生参数透传袋", OLIVE, False, 11),
    ], left_x + Inches(0.3), col_top + Inches(0.95),
       col_w - Inches(0.6), col_h - Inches(2.0),
       line_spacing=1.2)
    add_text(s,
             "入口层 _build_request 负责分拣:\n白名单进顶层,其余装进 provider_params。",
             left_x + Inches(0.3), col_top + Inches(3.7),
             col_w - Inches(0.6), Inches(0.8),
             font=SANS, size=11, color=DARK_WARM)

    # 右侧 · 透传模式字段
    add_text(s, "透传模式字段",
             right_x + Inches(0.3), col_top + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             font=SERIF, size=18, color=BRAND, bold=True)
    add_multi_text(s, [
        ("passthrough_mode", NEAR_BLACK, True, 13),
        ("    DISABLED / SAME_STYLE / CROSS_STYLE", OLIVE, False, 11),
        ("source_style", NEAR_BLACK, True, 13),
        ("    openai / anthropic / genai / bedrock", OLIVE, False, 11),
        ("raw_request", NEAR_BLACK, True, 13),
        ("    原生 SDK kwargs", OLIVE, False, 11),
        ("raw_method", NEAR_BLACK, True, 13),
        ("    chat.completions.create 等", OLIVE, False, 11),
    ], right_x + Inches(0.3), col_top + Inches(0.95),
       col_w - Inches(0.6), col_h - Inches(2.0),
       line_spacing=1.2)
    add_text(s,
             "Wrapped 入口不走分拣,所有参数原封不动\n塞进 raw_request,Hub 全程不拆包。",
             right_x + Inches(0.3), col_top + Inches(3.7),
             col_w - Inches(0.6), Inches(0.8),
             font=SANS, size=11, color=DARK_WARM)

    add_footer_note(s, "ADR-023 · fallback 阶段必须保留这 4 个字段,否则透传请求会退化成空的统一请求。")


# ═══════════════════════════════════════════════════════════
# slide 7 · CoreEngine 循环外/内
# ═══════════════════════════════════════════════════════════

def slide_engine(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "05 · Engine 节奏", 7)
    add_title(s, "CoreEngine 把工作拆成循环外一次,循环内每次重试都重新执行")

    col_top = Inches(2.0)
    col_h = Inches(4.6)
    col_w = Inches(6.0)
    left_x = Inches(0.6)
    right_x = Inches(6.93)

    # 左:循环外
    add_rect(s, left_x, col_top, col_w, col_h, fill=IVORY, border=BORDER, border_pt=0.5)
    add_text(s, "循环外 · 整个请求只跑一次",
             left_x + Inches(0.3), col_top + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             font=SERIF, size=16, color=OLIVE, bold=True)
    steps_left = [
        ("①", "_validate_request",
         "透传分支看 raw_*,统一分支看 messages"),
        ("②", "Plugin.before_request 链",
         "基础钩子,大多数插件在这步是 no-op"),
    ]
    for i, (num, title, body) in enumerate(steps_left):
        y = col_top + Inches(0.95 + i * 1.6)
        add_text(s, num, left_x + Inches(0.3), y, Inches(0.5), Inches(0.5),
                 font=SERIF, size=22, color=BRAND, bold=True)
        add_text(s, title, left_x + Inches(0.9), y + Inches(0.05),
                 col_w - Inches(1.2), Inches(0.4),
                 font=MONO, size=12, color=NEAR_BLACK, bold=True)
        add_text(s, body, left_x + Inches(0.9), y + Inches(0.5),
                 col_w - Inches(1.2), Inches(1.0),
                 font=SANS, size=11, color=DARK_WARM)

    # 右:循环内
    add_rect(s, right_x, col_top, col_w, col_h, fill=BRAND_TINT, border=BRAND, border_pt=1)
    add_text(s, "循环内 · 每次重试重跑",
             right_x + Inches(0.3), col_top + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             font=SERIF, size=16, color=BRAND, bold=True)
    steps_right = [
        ("③", "排除集 + 权重调整", "熔断器/限流器/自适应权重在此发力"),
        ("④", "Router.choose", "在剩余候选池中加权随机"),
        ("⑤", "_effective_request 合并参数", "request 显式 > endpoint > engine"),
        ("⑥", "Provider.invoke", "发真正的 HTTP 请求"),
        ("⑦", "成功 after / 失败 救援链", "回写滑窗 · 进 Retry → Failover → Fallback"),
    ]
    for i, (num, title, body) in enumerate(steps_right):
        y = col_top + Inches(0.85 + i * 0.72)
        add_text(s, num, right_x + Inches(0.3), y, Inches(0.4), Inches(0.4),
                 font=SERIF, size=15, color=BRAND, bold=True)
        add_text(s, title, right_x + Inches(0.75), y + Inches(0.02),
                 col_w - Inches(1.05), Inches(0.32),
                 font=MONO, size=10, color=NEAR_BLACK, bold=True)
        add_text(s, body, right_x + Inches(0.75), y + Inches(0.35),
                 col_w - Inches(1.05), Inches(0.3),
                 font=SANS, size=10, color=DARK_WARM)

    add_footer_note(s, "第 ③ 步是容错插件影响路由的真正入口 · Engine 在路由前主动询问插件。")


# ═══════════════════════════════════════════════════════════
# slide 8 · ★ Provider 三岔分派
# ═══════════════════════════════════════════════════════════

def slide_three_way(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "06 · 核心抽象 · ★", 8)
    add_title(s, "Provider 三岔分派:一段 if/else 决定了整条调用路径的形态")

    # 代码块
    code_x = Inches(0.6)
    code_top = Inches(2.0)
    code_w = Inches(12.13)
    code_h = Inches(1.85)
    add_rect(s, code_x, code_top, code_w, code_h,
             fill=RGBColor(0x2D, 0x2D, 0x2A), border=None,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)

    code_lines = [
        ("def invoke(self, request, decision) -> ModelResponse:", RGBColor(0xD8, 0xD8, 0xD2), False, 12),
        ("    if not request.is_passthrough:", RGBColor(0xD8, 0xD8, 0xD2), False, 12),
        ("        return self._invoke_unified(request, decision)        # 统一模式", RGBColor(0x9C, 0xCB, 0xEA), False, 12),
        ("    if request.source_style == self.api_style:", RGBColor(0xD8, 0xD8, 0xD2), False, 12),
        ("        return self._invoke_passthrough(request, decision)    # 同风格透传", RGBColor(0x9C, 0xCB, 0xEA), False, 12),
        ("    else:", RGBColor(0xD8, 0xD8, 0xD2), False, 12),
        ("        return self._invoke_adapted(request, decision)        # 跨风格透传", RGBColor(0x9C, 0xCB, 0xEA), False, 12),
    ]
    add_multi_text(s, code_lines,
                   code_x + Inches(0.4), code_top + Inches(0.15),
                   code_w - Inches(0.8), code_h - Inches(0.3),
                   font=MONO, size=12, color=WHITE, line_spacing=1.15)

    # 三栏表格
    tbl_top = Inches(4.15)
    tbl_h = Inches(2.6)
    headers = ["方法", "走它的场景", "转换次数", "谁实现"]
    rows = [
        ("_invoke_unified", "ChatModelHub / ModelHubClient", "2 次", "子类必须实现"),
        ("_invoke_passthrough", "Wrapped 路由到本家 endpoint", "0 次", "基类默认"),
        ("_invoke_adapted", "Wrapped 路由到异家 endpoint", "adapter 内 2 次", "基类默认"),
    ]
    col_widths = [Inches(3.4), Inches(4.5), Inches(1.8), Inches(2.43)]

    # 表头
    x = Inches(0.6)
    add_rect(s, x, tbl_top, Inches(12.13), Inches(0.5),
             fill=BRAND, border=None, shape=MSO_SHAPE.RECTANGLE)
    cur_x = x
    for i, h in enumerate(headers):
        add_text(s, h, cur_x + Inches(0.2), tbl_top + Inches(0.1),
                 col_widths[i] - Inches(0.2), Inches(0.3),
                 font=SANS, size=12, color=WHITE, bold=True)
        cur_x += col_widths[i]

    # 行
    for ri, row in enumerate(rows):
        y = tbl_top + Inches(0.5 + ri * 0.62)
        fill = IVORY if ri % 2 == 0 else PARCHMENT
        add_rect(s, x, y, Inches(12.13), Inches(0.62),
                 fill=fill, border=BORDER, border_pt=0.3, shape=MSO_SHAPE.RECTANGLE)
        cur_x = x
        for i, cell in enumerate(row):
            font = MONO if i == 0 else SANS
            color = NEAR_BLACK if i == 0 else DARK_WARM
            bold = i == 0
            add_text(s, cell, cur_x + Inches(0.2), y + Inches(0.18),
                     col_widths[i] - Inches(0.2), Inches(0.4),
                     font=font, size=12, color=color, bold=bold)
            cur_x += col_widths[i]

    add_footer_note(s, "子类只需要写 _invoke_unified · 透传与适配复用基类,新接一家供应商成本极低。")


# ═══════════════════════════════════════════════════════════
# slide 9 · 三层容错总览
# ═══════════════════════════════════════════════════════════

def slide_resilience_overview(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "07 · 容错全景", 9)
    add_title(s, "故障处理跑在两条独立轨道上,滑窗数据互相喂养形成闭环")

    # 左:救援链
    left_x = Inches(0.6)
    top1 = Inches(2.0)
    h1 = Inches(2.0)
    panel_w = Inches(6.0)
    add_rect(s, left_x, top1, panel_w, h1, fill=BRAND_TINT, border=BRAND, border_pt=1)
    add_text(s, "救援链 · 当前请求",
             left_x + Inches(0.3), top1 + Inches(0.18), panel_w - Inches(0.6), Inches(0.4),
             font=SERIF, size=14, color=BRAND, bold=True)

    sub_w = (panel_w - Inches(1.0)) / 3
    for i, (t, sub) in enumerate([
        ("退避重试", "同 endpoint"),
        ("Failover", "换 endpoint"),
        ("Fallback", "换备选模型"),
    ]):
        x = left_x + Inches(0.3) + (sub_w + Inches(0.2)) * i
        add_rect(s, x, top1 + Inches(0.75), sub_w, Inches(1.0),
                 fill=IVORY, border=BORDER, border_pt=0.5)
        add_text(s, t, x, top1 + Inches(0.85), sub_w, Inches(0.4),
                 font=SERIF, size=13, color=NEAR_BLACK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, sub, x, top1 + Inches(1.3), sub_w, Inches(0.4),
                 font=SANS, size=10, color=OLIVE, align=PP_ALIGN.CENTER)
        if i < 2:
            ax1 = x + sub_w
            ax2 = ax1 + Inches(0.2)
            ay = top1 + Inches(1.25)
            add_arrow(s, ax1, ay, ax2, ay, color=STONE, pt=1)

    # 右:健康管理
    right_x = Inches(6.93)
    add_rect(s, right_x, top1, panel_w, h1, fill=BRAND_TINT, border=BRAND, border_pt=1)
    add_text(s, "健康管理 · 影响未来请求",
             right_x + Inches(0.3), top1 + Inches(0.18), panel_w - Inches(0.6), Inches(0.4),
             font=SERIF, size=14, color=BRAND, bold=True)
    sub2_w = (panel_w - Inches(0.8)) / 2
    for i, (t, sub) in enumerate([
        ("自适应权重", "429 渐进降权"),
        ("熔断器", "5xx 完全摘除"),
    ]):
        x = right_x + Inches(0.3) + (sub2_w + Inches(0.2)) * i
        add_rect(s, x, top1 + Inches(0.75), sub2_w, Inches(1.0),
                 fill=IVORY, border=BORDER, border_pt=0.5)
        add_text(s, t, x, top1 + Inches(0.85), sub2_w, Inches(0.4),
                 font=SERIF, size=13, color=NEAR_BLACK, bold=True, align=PP_ALIGN.CENTER)
        add_text(s, sub, x, top1 + Inches(1.3), sub2_w, Inches(0.4),
                 font=SANS, size=10, color=OLIVE, align=PP_ALIGN.CENTER)

    # 闭环箭头
    add_arrow(s, left_x + panel_w, top1 + Inches(0.6),
              right_x, top1 + Inches(0.6), color=BRAND, pt=1.2)
    add_text(s, "不论成败 · 回写滑窗",
             left_x + panel_w, top1 + Inches(0.2),
             right_x - (left_x + panel_w), Inches(0.35),
             font=SANS, size=10, color=BRAND, align=PP_ALIGN.CENTER)
    add_arrow(s, right_x, top1 + Inches(1.55),
              left_x + panel_w, top1 + Inches(1.55), color=BRAND, pt=1.2)
    add_text(s, "下次路由 · 排除 + 降权",
             left_x + panel_w, top1 + Inches(1.65),
             right_x - (left_x + panel_w), Inches(0.35),
             font=SANS, size=10, color=BRAND, align=PP_ALIGN.CENTER)

    # 三种信号
    top2 = Inches(4.5)
    h2 = Inches(2.2)
    add_rect(s, left_x, top2, Inches(12.13), h2, fill=IVORY, border=BORDER, border_pt=0.5)
    add_text(s, "三种健康信号 · 按错误类型分工,互不重叠",
             left_x + Inches(0.3), top2 + Inches(0.2), Inches(11.5), Inches(0.4),
             font=SERIF, size=14, color=NEAR_BLACK, bold=True)
    sig_w = Inches(3.8)
    sig_top = top2 + Inches(0.8)
    sig_gap = Inches(0.2)
    for i, (sig, plug, how) in enumerate([
        ("429 限流", "自适应权重 + 限流器", "渐进降权 + 解析 Retry-After"),
        ("5xx 服务端错误", "熔断器", "CLOSED → OPEN → HALF_OPEN"),
        ("发送过快", "限流器", "令牌桶主动钳制"),
    ]):
        x = left_x + Inches(0.3) + (sig_w + sig_gap) * i
        add_text(s, sig, x, sig_top, sig_w, Inches(0.4),
                 font=SERIF, size=13, color=BRAND, bold=True)
        add_text(s, plug, x, sig_top + Inches(0.4), sig_w, Inches(0.4),
                 font=SANS, size=11, color=NEAR_BLACK)
        add_text(s, how, x, sig_top + Inches(0.8), sig_w, Inches(0.4),
                 font=SANS, size=10, color=OLIVE)


# ═══════════════════════════════════════════════════════════
# slide 10 · 救援链三级
# ═══════════════════════════════════════════════════════════

def slide_rescue_chain(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "08 · 救援链", 10)
    add_title(s, "三级递进容错 · 前一级失败才进入下一级")

    steps = [
        ("①", "退避重试", "当前 endpoint 返回 429",
         "按 Retry-After 退避\n同 endpoint 重试一次"),
        ("②", "Failover", "重试仍失败,或返回 5xx / 超时",
         "加入排除集\n从同模型其他 endpoint 选一个"),
        ("③", "Fallback", "同模型所有 endpoint 均不可用",
         "按 fallback_model 列表\n切到备选模型,独立享受 ①②"),
    ]
    card_w = Inches(3.93)
    gap = Inches(0.22)
    top = Inches(2.2)
    h = Inches(3.6)

    for i, (num, name, trigger, action) in enumerate(steps):
        x = Inches(0.6) + (card_w + gap) * i
        add_rect(s, x, top, card_w, h, fill=IVORY, border=BORDER, border_pt=0.6)
        add_text(s, num, x + Inches(0.3), top + Inches(0.25), Inches(1), Inches(0.8),
                 font=SERIF, size=44, color=BRAND, bold=True)
        add_text(s, name, x + Inches(0.3), top + Inches(1.15), card_w - Inches(0.6), Inches(0.5),
                 font=SERIF, size=20, color=NEAR_BLACK, bold=True)
        add_text(s, "触发", x + Inches(0.3), top + Inches(1.75), card_w - Inches(0.6), Inches(0.3),
                 font=MONO, size=10, color=STONE, bold=True)
        add_text(s, trigger, x + Inches(0.3), top + Inches(2.05), card_w - Inches(0.6), Inches(0.7),
                 font=SANS, size=12, color=DARK_WARM)
        add_text(s, "动作", x + Inches(0.3), top + Inches(2.5), card_w - Inches(0.6), Inches(0.3),
                 font=MONO, size=10, color=STONE, bold=True)
        add_text(s, action, x + Inches(0.3), top + Inches(2.8), card_w - Inches(0.6), Inches(0.8),
                 font=SANS, size=12, color=DARK_WARM)
        if i < 2:
            ax1 = x + card_w
            ax2 = ax1 + gap
            ay = top + h / 2
            add_arrow(s, ax1, ay, ax2, ay, color=BRAND, pt=1.5)

    add_footer_note(s, "ADR-023 · fallback 切换时必须保留 4 个透传字段,否则透传请求会退化。")


# ═══════════════════════════════════════════════════════════
# slide 11 · 熔断器
# ═══════════════════════════════════════════════════════════

def slide_circuit_breaker(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "09 · 熔断器", 11)
    add_title(s, "熔断器是二值的,处理坏了的场景,与降权天然分工")

    col_top = Inches(2.0)
    col_h = Inches(4.6)
    col_w = Inches(6.0)
    left_x = Inches(0.6)
    right_x = Inches(6.93)

    # 左:三状态机
    add_rect(s, left_x, col_top, col_w, col_h, fill=IVORY, border=BORDER, border_pt=0.5)
    add_text(s, "三状态机",
             left_x + Inches(0.3), col_top + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             font=SERIF, size=16, color=BRAND, bold=True)

    # 三个状态圆角矩形
    states = [("CLOSED", "正常参与路由", IVORY, BORDER),
              ("OPEN", "Router 完全跳过", BRAND_TINT, BRAND),
              ("HALF_OPEN", "放行少量探测", IVORY, BORDER)]
    sw = (col_w - Inches(0.8)) / 3
    sy = col_top + Inches(1.0)
    for i, (n, d, fill, bd) in enumerate(states):
        x = left_x + Inches(0.3) + (sw + Inches(0.1)) * i
        bw = 1 if bd == BRAND else 0.5
        add_rect(s, x, sy, sw, Inches(0.95), fill=fill, border=bd, border_pt=bw)
        add_text(s, n, x, sy + Inches(0.1), sw, Inches(0.35),
                 font=MONO, size=11, color=NEAR_BLACK, align=PP_ALIGN.CENTER, bold=True)
        add_text(s, d, x, sy + Inches(0.5), sw, Inches(0.4),
                 font=SANS, size=10, color=OLIVE, align=PP_ALIGN.CENTER)

    # 状态转移
    add_text(s, "状态转移",
             left_x + Inches(0.3), col_top + Inches(2.25), col_w - Inches(0.6), Inches(0.4),
             font=MONO, size=11, color=STONE, bold=True)
    add_multi_text(s, [
        ("CLOSED → OPEN", BRAND, True, 11),
        ("    滑窗失败率 > 50% 且调用 ≥ 5 次", DARK_WARM, False, 11),
        ("OPEN → HALF_OPEN", BRAND, True, 11),
        ("    冷却期 30s 到期", DARK_WARM, False, 11),
        ("HALF_OPEN → CLOSED", BRAND, True, 11),
        ("    连续成功 3 次", DARK_WARM, False, 11),
        ("HALF_OPEN → OPEN", BRAND, True, 11),
        ("    任意一次失败", DARK_WARM, False, 11),
    ], left_x + Inches(0.3), col_top + Inches(2.7),
       col_w - Inches(0.6), Inches(2.0), line_spacing=1.2)

    # 右:与自适应权重的分工
    add_rect(s, right_x, col_top, col_w, col_h, fill=IVORY, border=BORDER, border_pt=0.5)
    add_text(s, "与自适应权重的分工",
             right_x + Inches(0.3), col_top + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             font=SERIF, size=16, color=BRAND, bold=True)
    add_text(s, "熔断器是 0 / 1 二值,适合\"坏了\"",
             right_x + Inches(0.3), col_top + Inches(0.95), col_w - Inches(0.6), Inches(0.4),
             font=SANS, size=12, color=NEAR_BLACK)
    add_text(s, "自适应权重是连续渐进,适合\"太忙\"",
             right_x + Inches(0.3), col_top + Inches(1.35), col_w - Inches(0.6), Inches(0.4),
             font=SANS, size=12, color=NEAR_BLACK)
    add_line(s, right_x + Inches(0.3), col_top + Inches(1.85),
             right_x + col_w - Inches(0.3), col_top + Inches(1.85),
             color=BORDER, pt=0.5)

    # 错误类型隔离
    add_text(s, "错误类型隔离",
             right_x + Inches(0.3), col_top + Inches(2.0), col_w - Inches(0.6), Inches(0.4),
             font=MONO, size=11, color=STONE, bold=True)
    add_multi_text(s, [
        ("429 太忙", BRAND, True, 12),
        ("    自适应权重渐进降权,不触发熔断", DARK_WARM, False, 11),
        ("5xx 故障", BRAND, True, 12),
        ("    熔断器完全摘除,不影响降权", DARK_WARM, False, 11),
    ], right_x + Inches(0.3), col_top + Inches(2.45),
       col_w - Inches(0.6), Inches(1.5), line_spacing=1.3)

    add_text(s, "推荐配置 · exclude_429_from_circuit_breaker = true",
             right_x + Inches(0.3), col_top + Inches(4.0), col_w - Inches(0.6), Inches(0.5),
             font=MONO, size=10, color=BRAND, bold=True)


# ═══════════════════════════════════════════════════════════
# slide 12 · ★ 自适应权重 · 429 是太忙
# ═══════════════════════════════════════════════════════════

def slide_adaptive_weight(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "10 · 自适应权重 · ★", 12)
    add_title(s, "429 是太忙不是坏了 · 朴素的二值切换会把备份压垮")

    # 对比表
    tbl_top = Inches(2.0)
    headers = ["", "429 Too Many Requests", "5xx Server Error"]
    rows = [
        ("含义", "我没坏,只是请求太多了", "我真的出问题了"),
        ("正确应对", "少给它流量 · 降权", "暂时别用 · 熔断"),
        ("恢复预期", "减少请求后很快恢复", "需要等修复"),
    ]
    col_widths = [Inches(2.0), Inches(5.06), Inches(5.07)]
    x = Inches(0.6)
    add_rect(s, x, tbl_top, Inches(12.13), Inches(0.5),
             fill=BRAND, border=None, shape=MSO_SHAPE.RECTANGLE)
    cur_x = x
    for i, h in enumerate(headers):
        add_text(s, h, cur_x + Inches(0.2), tbl_top + Inches(0.1),
                 col_widths[i] - Inches(0.2), Inches(0.3),
                 font=SANS, size=12, color=WHITE, bold=True)
        cur_x += col_widths[i]
    for ri, row in enumerate(rows):
        y = tbl_top + Inches(0.5 + ri * 0.55)
        fill = IVORY if ri % 2 == 0 else PARCHMENT
        add_rect(s, x, y, Inches(12.13), Inches(0.55),
                 fill=fill, border=BORDER, border_pt=0.3, shape=MSO_SHAPE.RECTANGLE)
        cur_x = x
        for i, cell in enumerate(row):
            col = BRAND if i == 0 else DARK_WARM
            bold = i == 0
            add_text(s, cell, cur_x + Inches(0.2), y + Inches(0.15),
                     col_widths[i] - Inches(0.2), Inches(0.4),
                     font=SANS, size=12, color=col, bold=bold)
            cur_x += col_widths[i]

    # 公式块
    code_top = Inches(4.4)
    add_rect(s, Inches(0.6), code_top, Inches(12.13), Inches(1.2),
             fill=RGBColor(0x2D, 0x2D, 0x2A), border=None)
    add_multi_text(s, [
        ("error_rate  = 窗口内 429 次数 / 窗口内总调用次数", WHITE, False, 13),
        ("multiplier  = max(min_weight_ratio, 1.0 - error_rate × penalty_factor)", WHITE, False, 13),
        ("effective   = base_weight × multiplier", WHITE, False, 13),
    ], Inches(0.85), code_top + Inches(0.15),
       Inches(11.5), Inches(1.0),
       font=MONO, size=13, color=WHITE, line_spacing=1.2)

    add_text(s,
             "下限保护 · min_weight_ratio = 0.1 保留 10% 探测流量。这 10% 充当探测包,",
             Inches(0.6), Inches(5.85), Inches(12.13), Inches(0.4),
             font=SANS, size=12, color=DARK_WARM)
    add_text(s,
             "让系统不需要专门状态机就能感知配额恢复——业务流量本身就是探测。",
             Inches(0.6), Inches(6.2), Inches(12.13), Inches(0.4),
             font=SANS, size=12, color=DARK_WARM)

    add_footer_note(s, "熔断器靠 30s 冷却 + HALF_OPEN 主动探测;自适应权重靠这 10% 流量被动探测。")


# ═══════════════════════════════════════════════════════════
# slide 13 · 权重重分配 · 必要性
# ═══════════════════════════════════════════════════════════

def slide_redist_why(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "11 · 重分配 · 必要性 · ★", 13)
    add_title(s, "仅靠自降权,weight=1 的小 endpoint 永远到不了应急需要的占比")
    add_lead(s, "加权随机的本质 · 被选中概率 = weight_i / Σ weight。改变占比只有两条路。")

    # 两栏 · 分子/分母
    col_top = Inches(2.55)
    col_h = Inches(1.4)
    col_w = Inches(6.0)
    left_x = Inches(0.6)
    right_x = Inches(6.93)
    add_rect(s, left_x, col_top, col_w, col_h, fill=IVORY, border=BORDER, border_pt=0.5)
    add_text(s, "路径 1 · 分子变小", left_x + Inches(0.3), col_top + Inches(0.2),
             col_w - Inches(0.6), Inches(0.4),
             font=SERIF, size=14, color=BRAND, bold=True)
    add_text(s, "自降权 → 自己占比下降",
             left_x + Inches(0.3), col_top + Inches(0.65), col_w - Inches(0.6), Inches(0.4),
             font=SANS, size=12, color=DARK_WARM)
    add_text(s, "已实现 · 每个 endpoint 各自按 429 比例打 multiplier",
             left_x + Inches(0.3), col_top + Inches(1.0), col_w - Inches(0.6), Inches(0.4),
             font=SANS, size=11, color=OLIVE)

    add_rect(s, right_x, col_top, col_w, col_h, fill=BRAND_TINT, border=BRAND, border_pt=1)
    add_text(s, "路径 2 · 分母变小", right_x + Inches(0.3), col_top + Inches(0.2),
             col_w - Inches(0.6), Inches(0.4),
             font=SERIF, size=14, color=BRAND, bold=True)
    add_text(s, "别人降权 → 自己占比被动上升",
             right_x + Inches(0.3), col_top + Inches(0.65), col_w - Inches(0.6), Inches(0.4),
             font=SANS, size=12, color=DARK_WARM)
    add_text(s, "受 base weight 天花板限制 · weight=1 永远到不了高占比",
             right_x + Inches(0.3), col_top + Inches(1.0), col_w - Inches(0.6), Inches(0.4),
             font=SANS, size=11, color=OLIVE)

    # 数据示例
    ex_top = Inches(4.2)
    add_rect(s, Inches(0.6), ex_top, Inches(12.13), Inches(2.5),
             fill=IVORY, border=BORDER, border_pt=0.5)
    add_text(s, "配置 · gemini=66, gpt-5=34, gpt-4-1=1, o3=1",
             Inches(0.85), ex_top + Inches(0.2), Inches(11.5), Inches(0.4),
             font=MONO, size=12, color=NEAR_BLACK, bold=True)
    add_text(s, "gemini 与 gpt-5 都打到 multiplier=0.1 时:",
             Inches(0.85), ex_top + Inches(0.65), Inches(11.5), Inches(0.4),
             font=SANS, size=12, color=DARK_WARM)
    add_multi_text(s, [
        ("分母 = 6.6 + 3.4 + 1 + 1 = 12", NEAR_BLACK, False, 13),
        ("gpt-4-1 占比 = 1 / 12 = 8.3%", BRAND, True, 14),
        ("← 被 base=1 卡住的天花板,小 endpoint 接不住应急流量", OLIVE, False, 12),
    ], Inches(0.85), ex_top + Inches(1.15),
       Inches(11.5), Inches(1.3),
       font=MONO, color=DARK_WARM, line_spacing=1.5)

    add_footer_note(s, "要让小 endpoint 真正承接流量,必须显式把降权部分加到健康 endpoint 的分子上。")


# ═══════════════════════════════════════════════════════════
# slide 14 · 权重重分配 · 公式
# ═══════════════════════════════════════════════════════════

def slide_redist_formula(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "12 · 重分配 · 公式", 14)
    add_title(s, "把降权方释放的流量,按原权重比例分给完全健康的 endpoint")

    col_top = Inches(2.0)
    col_h = Inches(4.6)
    col_w = Inches(6.0)
    left_x = Inches(0.6)
    right_x = Inches(6.93)

    # 左 · 分组规则
    add_rect(s, left_x, col_top, col_w, col_h, fill=IVORY, border=BORDER, border_pt=0.5)
    add_text(s, "分组规则",
             left_x + Inches(0.3), col_top + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             font=SERIF, size=16, color=BRAND, bold=True)
    add_rect(s, left_x + Inches(0.3), col_top + Inches(0.9),
             col_w - Inches(0.6), Inches(1.4),
             fill=RGBColor(0x2D, 0x2D, 0x2A), border=None)
    add_multi_text(s, [
        ("D = { i : multiplier_i < 1.0 }  降权集", WHITE, False, 13),
        ("H = { i : multiplier_i = 1.0 }  健康集", WHITE, False, 13),
    ], left_x + Inches(0.5), col_top + Inches(1.05),
       col_w - Inches(0.9), Inches(1.1),
       font=MONO, size=13, color=WHITE, line_spacing=1.5)

    add_text(s, "为什么 H 严格用 = 1.0",
             left_x + Inches(0.3), col_top + Inches(2.5),
             col_w - Inches(0.6), Inches(0.4),
             font=SERIF, size=13, color=NEAR_BLACK, bold=True)
    add_text(s,
             "multiplier < 1.0 说明该 endpoint 自身已在 429,\n"
             "再把流量分给它只会加速把它也压到下限,\n"
             "引发连锁崩溃。只有完全健康的才安全接收。",
             left_x + Inches(0.3), col_top + Inches(2.95),
             col_w - Inches(0.6), Inches(1.5),
             font=SANS, size=12, color=DARK_WARM)

    # 右 · 公式
    add_rect(s, right_x, col_top, col_w, col_h, fill=IVORY, border=BORDER, border_pt=0.5)
    add_text(s, "重分配公式",
             right_x + Inches(0.3), col_top + Inches(0.25), col_w - Inches(0.6), Inches(0.5),
             font=SERIF, size=16, color=BRAND, bold=True)
    add_rect(s, right_x + Inches(0.3), col_top + Inches(0.9),
             col_w - Inches(0.6), Inches(2.5),
             fill=RGBColor(0x2D, 0x2D, 0x2A), border=None)
    add_multi_text(s, [
        ("Δ = Σ_{i∈D} base_i × (1 - multiplier_i)", WHITE, False, 12),
        ("", WHITE, False, 8),
        ("bonus_j = Δ × base_j / Σ_{k∈H} base_k", WHITE, False, 12),
        ("                                       (j∈H)", RGBColor(0x9C, 0xCB, 0xEA), False, 11),
        ("", WHITE, False, 8),
        ("降权方 · effective_i = base_i × multiplier_i", WHITE, False, 12),
        ("健康方 · effective_j = base_j + bonus_j", WHITE, False, 12),
    ], right_x + Inches(0.5), col_top + Inches(1.05),
       col_w - Inches(0.9), Inches(2.3),
       font=MONO, size=12, color=WHITE, line_spacing=1.2)

    add_multi_text(s, [
        ("性质 1 · 总权重守恒  Σ effective = Σ base", BRAND, True, 12),
        ("性质 2 · H 内相对比例严格保留", BRAND, True, 12),
    ], right_x + Inches(0.3), col_top + Inches(3.6),
       col_w - Inches(0.6), Inches(1.0),
       line_spacing=1.4)

    add_footer_note(s, "配置者写下 gpt-5:34, gpt-4-1:1 时已表达\"前者重要 34 倍\",应急时仍按这个比例分。")


# ═══════════════════════════════════════════════════════════
# slide 15 · 权重重分配 · 演练
# ═══════════════════════════════════════════════════════════

def slide_redist_walkthrough(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "13 · 重分配 · 演练", 15)
    add_title(s, "gemini=66 + gpt-5=34 同时 100% 429 · gpt-4-1 占比从 8.3% 升到 45.1%")

    # 代码块演算
    code_top = Inches(1.95)
    add_rect(s, Inches(0.6), code_top, Inches(12.13), Inches(2.1),
             fill=RGBColor(0x2D, 0x2D, 0x2A), border=None)
    add_multi_text(s, [
        ("Step 1 · 分组         D = { gemini, gpt-5 },  H = { gpt-4-1, o3 }", WHITE, False, 12),
        ("Step 2 · 减少总量     gemini 贡献 = 66 × 0.9 = 59.4", WHITE, False, 12),
        ("                      gpt-5  贡献 = 34 × 0.9 = 30.6", WHITE, False, 12),
        ("                      Δ = 90", RGBColor(0x9C, 0xCB, 0xEA), True, 12),
        ("Step 3 · bonus 分配   H 内 base 比例 1:1,各得 45", WHITE, False, 12),
        ("Step 4 · effective    gemini=6.6, gpt-5=3.4, gpt-4-1=46, o3=46", WHITE, False, 12),
        ("                      总和 = 102 = 原始总和  ← 权重守恒", RGBColor(0x9C, 0xCB, 0xEA), True, 12),
    ], Inches(0.85), code_top + Inches(0.15),
       Inches(11.5), Inches(1.85),
       font=MONO, size=12, color=WHITE, line_spacing=1.25)

    # 效果对照表
    tbl_top = Inches(4.3)
    headers = ["", "gemini", "gpt-5", "gpt-4-1", "o3"]
    rows = [
        ("无重分配 effective", "6.6", "3.4", "1", "1"),
        ("无重分配占比", "55.0%", "28.3%", "8.3%", "8.3%"),
        ("重分配后 effective", "6.6", "3.4", "46", "46"),
        ("重分配后占比", "6.5%", "3.3%", "45.1%", "45.1%"),
    ]
    col_widths = [Inches(3.43), Inches(2.17), Inches(2.17), Inches(2.18), Inches(2.18)]
    x = Inches(0.6)
    add_rect(s, x, tbl_top, Inches(12.13), Inches(0.45),
             fill=BRAND, border=None, shape=MSO_SHAPE.RECTANGLE)
    cur_x = x
    for i, h in enumerate(headers):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
        add_text(s, h, cur_x + Inches(0.15), tbl_top + Inches(0.08),
                 col_widths[i] - Inches(0.2), Inches(0.3),
                 font=SANS, size=12, color=WHITE, bold=True, align=align)
        cur_x += col_widths[i]
    for ri, row in enumerate(rows):
        y = tbl_top + Inches(0.45 + ri * 0.45)
        fill = IVORY if ri % 2 == 0 else PARCHMENT
        bold_row = ri >= 2
        add_rect(s, x, y, Inches(12.13), Inches(0.45),
                 fill=fill, border=BORDER, border_pt=0.3, shape=MSO_SHAPE.RECTANGLE)
        cur_x = x
        for i, cell in enumerate(row):
            col = BRAND if bold_row and i >= 3 else NEAR_BLACK
            align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
            font = MONO if i > 0 else SANS
            add_text(s, cell, cur_x + Inches(0.15), y + Inches(0.1),
                     col_widths[i] - Inches(0.2), Inches(0.35),
                     font=font, size=12, color=col, bold=(bold_row and i >= 3),
                     align=align)
            cur_x += col_widths[i]


# ═══════════════════════════════════════════════════════════
# slide 16 · recovery_rate
# ═══════════════════════════════════════════════════════════

def slide_recovery_rate(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "14 · recovery_rate · ★", 16)
    add_title(s, "恢复必须限速,否则瞬间涌回会让刚恢复的 endpoint 再次 429")
    add_lead(s, "重分配让恢复时刻同时双向变化——主力涨、备份跌,冲击比纯自降权更强。")

    # 公式块
    code_top = Inches(2.55)
    add_rect(s, Inches(0.6), code_top, Inches(12.13), Inches(1.6),
             fill=RGBColor(0x2D, 0x2D, 0x2A), border=None)
    add_multi_text(s, [
        ("核心规则 · 降权立即生效,恢复线性爬升", RGBColor(0x9C, 0xCB, 0xEA), True, 12),
        ("", WHITE, False, 6),
        ("if target < prev_multiplier:", WHITE, False, 12),
        ("    multiplier = target                            # 下行 · 立即生效", WHITE, False, 12),
        ("else:", WHITE, False, 12),
        ("    cap = start_multiplier + recovery_rate × elapsed", WHITE, False, 12),
        ("    multiplier = min(target, cap)                  # 上行 · 限速爬升", WHITE, False, 12),
    ], Inches(0.85), code_top + Inches(0.15),
       Inches(11.5), Inches(1.4),
       font=MONO, size=12, color=WHITE, line_spacing=1.2)

    # 恢复时间线
    tbl_top = Inches(4.45)
    headers = ["时刻", "multiplier", "流量份额分布"]
    rows = [
        ("t = 180s 之前", "0.10", "gemini 6.5%  ·  gpt-4-1 45.1%"),
        ("t = 180s 瞬间恢复 · 无 recovery_rate", "1.00", "gemini 64.7%  ·  gpt-4-1 1.0%  ·  ×11 涌回再次 429"),
        ("t = 183s · recovery_rate = 0.05", "0.25", "gemini 16.2%  ·  gpt-4-1 37.7%"),
        ("t = 198s · 18s 后爬满", "1.00", "gemini 64.7%  ·  gpt-4-1 1.0%  ·  平滑过渡"),
    ]
    col_widths = [Inches(4.5), Inches(1.6), Inches(6.03)]
    x = Inches(0.6)
    add_rect(s, x, tbl_top, Inches(12.13), Inches(0.4),
             fill=BRAND, border=None, shape=MSO_SHAPE.RECTANGLE)
    cur_x = x
    for i, h in enumerate(headers):
        add_text(s, h, cur_x + Inches(0.15), tbl_top + Inches(0.07),
                 col_widths[i] - Inches(0.2), Inches(0.3),
                 font=SANS, size=11, color=WHITE, bold=True)
        cur_x += col_widths[i]
    for ri, row in enumerate(rows):
        y = tbl_top + Inches(0.4 + ri * 0.42)
        fill = IVORY if ri % 2 == 0 else PARCHMENT
        bad = ri == 1
        add_rect(s, x, y, Inches(12.13), Inches(0.42),
                 fill=fill, border=BORDER, border_pt=0.3, shape=MSO_SHAPE.RECTANGLE)
        cur_x = x
        for i, cell in enumerate(row):
            if bad and i >= 1:
                col = RGBColor(0xb0, 0x40, 0x40)
            else:
                col = BRAND if i == 1 else DARK_WARM
            font = MONO if i >= 1 else SANS
            add_text(s, cell, cur_x + Inches(0.15), y + Inches(0.09),
                     col_widths[i] - Inches(0.2), Inches(0.3),
                     font=font, size=11, color=col, bold=(i == 1))
            cur_x += col_widths[i]


# ═══════════════════════════════════════════════════════════
# slide 17 · 完整调用旅程
# ═══════════════════════════════════════════════════════════

def slide_journey(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "15 · 完整旅程", 17)
    add_title(s, "一次 WrappedOpenAI 调用遇 429 · 自动跨厂商 fallback 到 Anthropic")

    steps = [
        ("1", "Wrapper 拦截",     "打包成 ModelRequest(SAME_STYLE,\nsource_style=openai, raw_request=...)"),
        ("2", "Transport",        "DirectProviderTransport.invoke\n进入 engine"),
        ("3", "Engine 路由",      "排除 OPEN 熔断与限流 endpoint\n降权后选中 endpoint A · OpenAI"),
        ("4", "Provider 三岔",    "source 与 api 都是 openai\n→ _invoke_passthrough · 0 转换"),
        ("5", "A 返回 429",       "回写 CB / AW 滑窗 · 退避重试仍失败\n→ Failover 到 endpoint B · Anthropic"),
        ("6", "进 Anthropic",     "source=openai ≠ api=anthropic\n→ _invoke_adapted"),
        ("7", "Adapter 翻译",     "openai_to_anthropic 翻译含图片\ncontent block,调 Anthropic SDK"),
        ("8", "响应回业务",       "adapt_response 反翻成 OpenAI 格式\nChatCompletion,Wrapper 拆包还给业务"),
    ]
    # 4×2 网格
    card_w = Inches(2.93)
    card_h = Inches(1.95)
    col_gap = Inches(0.16)
    row_gap = Inches(0.15)
    start_x = Inches(0.6)
    start_y = Inches(1.95)

    for idx, (num, title, body) in enumerate(steps):
        col = idx % 4
        row = idx // 4
        x = start_x + (card_w + col_gap) * col
        y = start_y + (card_h + row_gap) * row
        add_rect(s, x, y, card_w, card_h, fill=IVORY, border=BORDER, border_pt=0.5)
        # 步骤序号圆
        add_rect(s, x + Inches(0.2), y + Inches(0.2), Inches(0.45), Inches(0.45),
                 fill=BRAND, border=None, shape=MSO_SHAPE.OVAL)
        add_text(s, num, x + Inches(0.2), y + Inches(0.2), Inches(0.45), Inches(0.45),
                 font=SERIF, size=14, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, vanchor=MSO_ANCHOR.MIDDLE)
        add_text(s, title, x + Inches(0.75), y + Inches(0.25), card_w - Inches(0.95), Inches(0.5),
                 font=SERIF, size=14, color=NEAR_BLACK, bold=True)
        add_text(s, body, x + Inches(0.2), y + Inches(0.85), card_w - Inches(0.4), Inches(1.05),
                 font=SANS, size=10, color=DARK_WARM)

    add_footer_note(s, "全程业务方只写了一行 wrap_openai · 跨厂商、跨风格、跨容错完全无感。")


# ═══════════════════════════════════════════════════════════
# slide 18 · 能力复用矩阵
# ═══════════════════════════════════════════════════════════

def slide_capability_matrix(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "16 · 能力矩阵", 18)
    add_title(s, "三种模式共享同一套基础设施,任何模式都享受完整能力栈")

    headers = ["能力", "DISABLED · 统一", "SAME_STYLE · 同风格", "CROSS_STYLE · 跨风格"]
    rows = [
        ("路由决策 · 会话粘性", "✓", "✓", "✓"),
        ("重试 · 429 退避", "✓", "✓", "✓"),
        ("熔断 · 限流", "✓", "✓", "✓"),
        ("Fallback 救援", "✓", "✓", "✓"),
        ("before_request 插件", "读 messages", "读 raw_request", "读 raw_request"),
        ("after_response 插件", "读 choices", "读 raw_response", "读 raw_response"),
        ("数据转换次数", "2 次", "0 次", "2 次 · adapter 内"),
    ]
    col_widths = [Inches(3.43), Inches(2.9), Inches(2.9), Inches(2.9)]
    x = Inches(0.6)
    tbl_top = Inches(2.05)

    add_rect(s, x, tbl_top, Inches(12.13), Inches(0.5),
             fill=BRAND, border=None, shape=MSO_SHAPE.RECTANGLE)
    cur_x = x
    for i, h in enumerate(headers):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER
        add_text(s, h, cur_x + Inches(0.15), tbl_top + Inches(0.12),
                 col_widths[i] - Inches(0.2), Inches(0.3),
                 font=SANS, size=12, color=WHITE, bold=True, align=align)
        cur_x += col_widths[i]
    for ri, row in enumerate(rows):
        y = tbl_top + Inches(0.5 + ri * 0.55)
        fill = IVORY if ri % 2 == 0 else PARCHMENT
        add_rect(s, x, y, Inches(12.13), Inches(0.55),
                 fill=fill, border=BORDER, border_pt=0.3, shape=MSO_SHAPE.RECTANGLE)
        cur_x = x
        for i, cell in enumerate(row):
            if i == 0:
                col = NEAR_BLACK
                bold = True
                font = SANS
                align = PP_ALIGN.LEFT
            elif cell == "✓":
                col = BRAND
                bold = True
                font = SANS
                align = PP_ALIGN.CENTER
            else:
                col = DARK_WARM
                bold = False
                font = SANS
                align = PP_ALIGN.CENTER
            add_text(s, cell, cur_x + Inches(0.15), y + Inches(0.17),
                     col_widths[i] - Inches(0.2), Inches(0.35),
                     font=font, size=12, color=col, bold=bold, align=align)
            cur_x += col_widths[i]

    add_footer_note(s, "能力栈对三种模式无差异,业务方可按场景自由切换,基础设施代价不重复支付。")


# ═══════════════════════════════════════════════════════════
# slide 19 · 设计哲学
# ═══════════════════════════════════════════════════════════

def slide_principles(prs):
    s = blank_slide(prs)
    add_eyebrow(s, "17 · 设计哲学", 19)
    add_title(s, "四条设计原则把扩 Provider 的代价压到最低")

    principles = [
        ("A", "单一入口",
         "整个 Model Hub 对外只有一个签名\ninvoke(request) → response。\n基础设施只需实现一次,所有模式复用。"),
        ("B", "has-a 组合",
         "入口 → Transport → Engine → Registry\n→ Provider → SDK,层层组合不嵌套。"),
        ("C", "翻译是 provider 内政",
         "CROSS_STYLE 翻译完全在 provider 内做。\nEngine 不知道 adapter 存在,\nadapter 独立扩展不污染 engine。"),
        ("D", "数据结构扩展,而非入口扩展",
         "扩 ModelRequest 字段而非加新方法。\n代价是每 provider 加 3 个内部方法,\n收益是单入口与插件零双写。"),
    ]
    card_w = Inches(6.0)
    card_h = Inches(2.25)
    gap_x = Inches(0.13)
    gap_y = Inches(0.15)
    start_x = Inches(0.6)
    start_y = Inches(1.95)

    for idx, (letter, name, body) in enumerate(principles):
        col = idx % 2
        row = idx // 2
        x = start_x + (card_w + gap_x) * col
        y = start_y + (card_h + gap_y) * row
        add_rect(s, x, y, card_w, card_h, fill=IVORY, border=BORDER, border_pt=0.5)
        # 大字母
        add_text(s, letter, x + Inches(0.3), y + Inches(0.2), Inches(1.0), Inches(1.0),
                 font=SERIF, size=44, color=BRAND, bold=True)
        # 标题
        add_text(s, name, x + Inches(1.4), y + Inches(0.3),
                 card_w - Inches(1.6), Inches(0.5),
                 font=SERIF, size=18, color=NEAR_BLACK, bold=True)
        add_text(s, body, x + Inches(1.4), y + Inches(0.85),
                 card_w - Inches(1.6), Inches(1.3),
                 font=SANS, size=12, color=DARK_WARM)


# ═══════════════════════════════════════════════════════════
# slide 20 · 收尾
# ═══════════════════════════════════════════════════════════

def slide_ending(prs):
    s = blank_slide(prs)
    add_text(s, "谢谢倾听",
             Inches(0.6), Inches(2.8), Inches(12.13), Inches(1.2),
             font=SERIF, size=52, color=NEAR_BLACK,
             align=PP_ALIGN.CENTER)
    add_line(s, Inches(6.17), Inches(4.3), Inches(7.17), Inches(4.3),
             color=BRAND, pt=1.5)
    add_text(s, "Plaud Model Hub · 把多供应商 LLM 调度的共性下沉成基础设施",
             Inches(0.6), Inches(4.5), Inches(12.13), Inches(0.5),
             font=SANS, size=16, color=OLIVE, align=PP_ALIGN.CENTER)
    add_text(s, "详细设计 · docs/arch-v2.md   ·   ADR-006 / 023 / 024 / 025",
             Inches(0.6), Inches(6.6), Inches(12.13), Inches(0.4),
             font=MONO, size=11, color=STONE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# main
# ═══════════════════════════════════════════════════════════

def main():
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--out", default="PlaudModelHub技术分享.pptx",
                        help="输出 .pptx 路径")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_toc(prs)
    slide_pain(prs)
    slide_modes(prs)
    slide_architecture(prs)
    slide_modelrequest(prs)
    slide_engine(prs)
    slide_three_way(prs)
    slide_resilience_overview(prs)
    slide_rescue_chain(prs)
    slide_circuit_breaker(prs)
    slide_adaptive_weight(prs)
    slide_redist_why(prs)
    slide_redist_formula(prs)
    slide_redist_walkthrough(prs)
    slide_recovery_rate(prs)
    slide_journey(prs)
    slide_capability_matrix(prs)
    slide_principles(prs)
    slide_ending(prs)

    prs.save(args.out)
    print(f"OK · 共 {len(prs.slides)} 张幻灯片 · 已保存到 {args.out}")


if __name__ == "__main__":
    main()
